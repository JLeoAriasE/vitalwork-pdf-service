#!/usr/bin/env python3
"""
generar_pptx_psicosocial.py
Genera informe PPTX de evaluación psicosocial usando plantilla.
Soporta: Ministerio de Trabajo Ecuador y FPsico (NTP 926)
"""
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData
from pptx.oxml.ns import qn
from lxml import etree
import io, os

VERDE    = RGBColor(0x1E,0x84,0x49)
AMARILLO = RGBColor(0xFF,0xC0,0x00)
ROJO     = RGBColor(0xCC,0x11,0x11)
COL_NV   = {'Bajo':VERDE, 'Medio':AMARILLO, 'Alto':ROJO}
NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), 'plantilla_psicosocial.pptx')

def reemplazar_texto(shape, reemplazos):
    if not shape.has_text_frame: return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            for old, new in reemplazos.items():
                if old in run.text:
                    run.text = run.text.replace(old, new)

def actualizar_pie(slide, cats, vals):
    charts = [s for s in slide.shapes if s.shape_type == 3]
    if not charts: return
    cd = ChartData(); cd.categories = cats; cd.add_series('', vals)
    charts[0].chart.replace_data(cd)

def limpiar_parrafo(para):
    p = para._p
    for child in list(p):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('r','br','fld'): p.remove(child)

def escribir_texto(para, texto, sz=20, bold=True, color='002060', spc=1600):
    limpiar_parrafo(para)
    p = para._p
    pPr = p.find(f'{{{NS}}}pPr')
    if pPr is None:
        pPr = etree.Element(f'{{{NS}}}pPr')
        p.insert(0, pPr)
    for child in list(pPr):
        if child.tag.split('}')[-1] in ('spcBef','spcAft'): pPr.remove(child)
    if spc > 0:
        sb = etree.SubElement(pPr, f'{{{NS}}}spcBef')
        sp = etree.SubElement(sb, f'{{{NS}}}spcPts'); sp.set('val', str(spc))
    r = etree.SubElement(p, f'{{{NS}}}r')
    rPr = etree.SubElement(r, f'{{{NS}}}rPr')
    rPr.set('sz', str(int(sz*100))); rPr.set('b', '1' if bold else '0'); rPr.set('dirty', '0')
    sf = etree.SubElement(rPr, f'{{{NS}}}solidFill')
    clr = etree.SubElement(sf, f'{{{NS}}}srgbClr'); clr.set('val', color)
    t = etree.SubElement(r, f'{{{NS}}}t'); t.text = texto
    endPr = p.find(f'{{{NS}}}endParaRPr')
    if endPr is not None: p.remove(endPr); p.append(endPr)

def escribir_circulo(para, nv, label, pts, sz=28, spc=400):
    limpiar_parrafo(para)
    p = para._p
    pPr = p.find(f'{{{NS}}}pPr')
    if pPr is None:
        pPr = etree.Element(f'{{{NS}}}pPr'); p.insert(0, pPr)
    for child in list(pPr):
        if child.tag.split('}')[-1] in ('spcBef','spcAft'): pPr.remove(child)
    if spc > 0:
        sb = etree.SubElement(pPr, f'{{{NS}}}spcBef')
        sp = etree.SubElement(sb, f'{{{NS}}}spcPts'); sp.set('val', str(spc))
    col = COL_NV[nv]; hex_col = f'{col[0]:02X}{col[1]:02X}{col[2]:02X}'
    sz100 = int(sz*100)
    r1 = etree.SubElement(p, f'{{{NS}}}r')
    rPr1 = etree.SubElement(r1, f'{{{NS}}}rPr')
    rPr1.set('sz', str(sz100)); rPr1.set('b', '1'); rPr1.set('dirty', '0')
    sf1 = etree.SubElement(rPr1, f'{{{NS}}}solidFill')
    clr1 = etree.SubElement(sf1, f'{{{NS}}}srgbClr'); clr1.set('val', hex_col)
    t1 = etree.SubElement(r1, f'{{{NS}}}t'); t1.text = u'\u25cf  '
    r2 = etree.SubElement(p, f'{{{NS}}}r')
    rPr2 = etree.SubElement(r2, f'{{{NS}}}rPr')
    rPr2.set('sz', str(sz100)); rPr2.set('b', '1'); rPr2.set('dirty', '0')
    sf2 = etree.SubElement(rPr2, f'{{{NS}}}solidFill')
    clr2 = etree.SubElement(sf2, f'{{{NS}}}srgbClr'); clr2.set('val', '002060')
    t2 = etree.SubElement(r2, f'{{{NS}}}t'); t2.text = f'{label}: {pts} pts'
    endPr = p.find(f'{{{NS}}}endParaRPr')
    if endPr is not None: p.remove(endPr); p.append(endPr)

def reemplazar_shape_completo(shape, texto, sz=13, bold=False):
    if not shape.has_text_frame: return
    para = shape.text_frame.paragraphs[0]
    limpiar_parrafo(para)
    r = etree.SubElement(para._p, f'{{{NS}}}r')
    rPr = etree.SubElement(r, f'{{{NS}}}rPr')
    rPr.set('sz', str(int(sz*100))); rPr.set('b', '1' if bold else '0'); rPr.set('dirty', '0')
    sf = etree.SubElement(rPr, f'{{{NS}}}solidFill')
    clr = etree.SubElement(sf, f'{{{NS}}}srgbClr'); clr.set('val', '002060')
    t = etree.SubElement(r, f'{{{NS}}}t'); t.text = texto

def llenar_textbox(tf, items, fn_escribir):
    paras = tf.paragraphs
    for k, item in enumerate(items):
        if k < len(paras): fn_escribir(paras[k], item)
    for k in range(len(items), len(paras)): limpiar_parrafo(paras[k])

def eliminar_vacios(slide, shape_name):
    for shape in slide.shapes:
        if shape.name == shape_name and shape.has_text_frame:
            txBody = shape.text_frame._txBody
            paras = txBody.findall(f'{{{NS}}}p')
            ultimo = -1
            for j, p in enumerate(paras):
                txt = ''.join(t.text or '' for t in p.findall(f'.//{{{NS}}}t'))
                if txt.strip(): ultimo = j
            for j in range(len(paras)-1, ultimo, -1):
                txBody.remove(paras[j])
            break

def corregir_endpararpr(slide, shape_name):
    for shape in slide.shapes:
        if shape.name == shape_name and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                p = para._p
                endPr = p.find(f'{{{NS}}}endParaRPr')
                if endPr is not None:
                    p.remove(endPr); p.append(endPr)
            break

def generar_informe_psicosocial(data):
    """
    data: {
        empresa, ruc, actividad, representante, fecha,
        tipo_test ('Ministerio de Trabajo Ecuador' o 'FPsico NTP 926'),
        totalN, masc, fem,
        edades: {'18-24 años':N, ...},
        dims: [(label, pts, nivel), ...],
        ec: {nombre, direccion, wa_phone}
    }
    Retorna: bytes del PPTX
    """
    empresa      = data.get('empresa', '—')
    ruc          = data.get('ruc', '—')
    actividad    = data.get('actividad', '—')
    fecha        = data.get('fecha', '—')
    tipo_test    = data.get('tipo_test', 'Ministerio de Trabajo Ecuador')
    totalN       = data.get('totalN', 0)
    masc         = data.get('masc', 0)
    fem          = data.get('fem', 0)
    edades       = data.get('edades', {})
    dims         = data.get('dims', [])
    totalGlobal  = sum(d[1] for d in dims)
    nivel_g      = 'Bajo' if totalGlobal >= 175 else 'Medio' if totalGlobal >= 117 else 'Alto'
    col_map      = {'Bajo': VERDE, 'Medio': AMARILLO, 'Alto': ROJO}
    criticas     = [d for d in dims if d[2] in ['Alto', 'Medio']]

    prs = Presentation(TEMPLATE_PATH)

    # SLIDE 1 - Portada
    s1 = prs.slides[0]
    for shape in s1.shapes:
        if shape.name == 'CuadroTexto 12':
            reemplazar_shape_completo(shape, empresa, sz=36, bold=True)
        elif shape.name == 'TextBox 6':
            reemplazar_shape_completo(shape, f'TEST: {tipo_test}', sz=13)
        else:
            reemplazar_texto(shape, {
                'RUC: 0791841237001': f'RUC: {ruc}',
                'G4630.31': actividad,
                'Universo: 25 trabajadores': f'Universo: {totalN} trabajadores',
                '2026-03-18': fecha,
            })

    # SLIDE 2 - Género
    s2 = prs.slides[1]
    for shape in s2.shapes:
        reemplazar_texto(shape, {
            'Total trabajadores: 25': f'Total trabajadores: {totalN}',
            'Masculino: 18': f'Masculino: {masc}',
            'Femenino: 7': f'Femenino: {fem}',
        })
    actualizar_pie(s2, ['Masculino', 'Femenino'], [masc/max(totalN,1), fem/max(totalN,1)])

    # SLIDE 3 - Edades
    s3 = prs.slides[2]
    for shape in s3.shapes:
        reemplazar_texto(shape, {
            'Total trabajadores: 25': f'Total trabajadores: {totalN}',
            '18-24 años: 3':  f'18-24 años: {edades.get("18-24 años", 0)}',
            '25-34 años: 10': f'25-34 años: {edades.get("25-34 años", 0)}',
            '35-44 años: 8':  f'35-44 años: {edades.get("35-44 años", 0)}',
            '45-54 años: 3':  f'45-54 años: {edades.get("45-54 años", 0)}',
            '55 años en adelante: 1': f'55 años en adelante: {edades.get("55 años en adelante", 0)}',
        })
    actualizar_pie(s3, list(edades.keys()), [v/max(totalN,1) for v in edades.values()])

    # SLIDE 4 - Resultado Global
    s4 = prs.slides[3]
    for shape in s4.shapes:
        reemplazar_texto(shape, {
            'Puntaje Total: 151': f'Puntaje Total: {totalGlobal}',
            'Universo: 25 trabajadores': f'Universo: {totalN} trabajadores',
        })
        if shape.name == 'Rectángulo 8':
            shape.fill.solid(); shape.fill.fore_color.rgb = col_map[nivel_g]
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if 'RIESGO' in run.text:
                            run.text = f'RIESGO\n{nivel_g.upper()}'

    # SLIDE 5 - Resultados por Dimensión
    s5 = prs.slides[4]
    for shape in s5.shapes:
        if shape.name == 'TextBox 4' and shape.has_text_frame:
            shape.top = Inches(1.35); shape.height = Inches(5.4)
            llenar_textbox(shape.text_frame, dims,
                lambda p, d: escribir_circulo(p, d[2], d[0], d[1], sz=28, spc=400))
            corregir_endpararpr(s5, 'TextBox 4')
            eliminar_vacios(s5, 'TextBox 4')
            break

    # SLIDE 6 - Gráfico barras
    s6 = prs.slides[5]
    charts = [s for s in s6.shapes if s.shape_type == 3]
    if charts:
        cd = ChartData()
        cd.categories = [d[0] for d in dims]
        cd.add_series('Puntaje', [d[1] for d in dims])
        charts[0].chart.replace_data(cd)

    # SLIDE 7 - Dimensiones críticas
    s7 = prs.slides[6]
    for shape in s7.shapes:
        if shape.name == 'TextBox 4' and shape.has_text_frame:
            shape.top = Inches(1.4); shape.height = Inches(5.3)
            n = len(criticas)
            spc = {1:2400,2:1800,3:1200,4:800,5:600,6:400,7:400}.get(n, 400)
            sz  = {1:40,  2:36,  3:32,  4:30, 5:28, 6:26, 7:24 }.get(n, 24)
            llenar_textbox(shape.text_frame, criticas,
                lambda p, d: escribir_circulo(p, d[2], d[0], d[1], sz=sz, spc=spc))
            corregir_endpararpr(s7, 'TextBox 4')
            eliminar_vacios(s7, 'TextBox 4')
            break

    # SLIDE 8 - Recomendaciones
    s8 = prs.slides[7]
    altos = [d[0] for d in dims if d[2] == 'Alto']
    recs = [
        'Promover espacios de comunicación y retroalimentación entre líderes y equipos.',
        'Implementar pausas activas y programas de bienestar laboral.',
        'Fortalecer el plan de capacitación y desarrollo de competencias.',
        'Establecer mecanismos de seguimiento periódico de factores psicosociales.',
        'Socializar los resultados con todos los trabajadores y áreas.',
        'Implementar el Programa de Riesgos Psicosociales según normativa vigente.',
    ]
    if altos:
        recs.insert(0, f'Atención inmediata en: {", ".join(altos)}.')
    for shape in s8.shapes:
        if shape.name == 'TextBox 4' and shape.has_text_frame:
            shape.top = Inches(1.4); shape.height = Inches(5.3)
            llenar_textbox(shape.text_frame, recs,
                lambda p, r: escribir_texto(p, r, sz=26, spc=1600))
            break

    # SLIDE 9 - Cierre
    s9 = prs.slides[8]
    for shape in s9.shapes:
        reemplazar_texto(shape, {'2026-03-18': fecha})

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
